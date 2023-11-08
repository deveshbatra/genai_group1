import numpy as np
import pandas as pd
import openai
from pptx import Presentation
from pptx.util import Cm, Pt
from docx import Document
from pypdf import PdfReader
import json

# Replace 'your_api_key_here' with your actual OpenAI API key
openai.api_key = ''

from openai import OpenAI
client = OpenAI(api_key = '')
# Function to reword text using GPT-4
#def reword_text_with_gpt4(text):
#    try:
#        # Use the OpenAI API to get a response
#        response = openai.chat.completions.create(
#            model="gpt-3.5-turbo-1106",  # Replace with the appropriate GPT-4 model when available
#            messages=[
#          {"role": "system", "content": "Reword the following text to be clear and concise:\n\n" + text}
#                ]
#            #prompt="Reword the following text to be clear and concise:\n\n" + text,
#            #max_tokens=60  # Adjust max tokens as needed
#        )
#        return response.choices[0].text.strip()
#    except Exception as e:  # Catch a general exception
#        print(f"An error occurred: {e}")
#        return text  # Return the original text if an error occurs

def reword_text_with_gpt4(
    text: str, 
    audience_type: str, 
    shape_type: str) -> str:
    if len(text) == 0:
        return text
    try:
        if shape_type =="Title 1":
            title_prompt = " This is a title of a slide so keep it to less than 8 words."
        else:
            title_prompt =""
        # Use the OpenAI API to get a response
        response = openai.chat.completions.create(
            model="gpt-4-0613",  # Replace with the appropriate model
            messages=[
                {
                    "role": "system", 
                    "content": (
                        "You are an assistant that rewords sentences to be clear and concise. Your output will be no longer than the input in length." +
                        " I am presenting to " + audience_type + ", so make it suitable for this audience." + title_prompt

                                )
                    },
                {"role": "user", "content": text}
            ]
        )
        # Assuming the last message in the list will be the assistant's response
        return str(response.choices[0].message.content)#.text.strip()
    except Exception as e:  # Catch a general exception
        print(f"An error occurred: {e}")
        return text  # Return the original text if an error occurs

def create_exec_summary(
    text: str, 
    audience_type: str) -> str:
    try:
        # Use the OpenAI API to get a response
        response = openai.chat.completions.create(
            model="gpt-4-0613",  # Replace with the appropriate model
            messages=[
                {
                    "role": "system", 
                    "content": (
                        "You are an assistant that creates executive summaries." +
                        " I am presenting to " + audience_type + ", so make it suitable for this audience. Produce a simple five bullet point summary"

                                )
                    },
                {"role": "user", "content": text}
            ]
        )
        # Assuming the last message in the list will be the assistant's response
        return str(response.choices[0].message.content)#.text.strip()
    except Exception as e:  # Catch a general exception
        print(f"An error occurred: {e}")
        return text  # Return the original text if an error occurs

def create_ppt_feedback(
    text: str, 
    audience_type: str) -> str:
    try:
        # Use the OpenAI API to get a response
        response = openai.chat.completions.create(
            model="gpt-4-0613",  # Replace with the appropriate model
            messages=[
                {
                    "role": "system", 
                    "content": (
                        "You are an assistant that provides feedback for powerpoint presentations." +
                        " I am presenting to " + audience_type + ", so make it suitable for this audience." +
                       " Produce feedback on the following powerpoint contents and tell me how to improve it."

                                )
                    },
                {"role": "user", "content": text}
            ]
        )
        # Assuming the last message in the list will be the assistant's response
        return str(response.choices[0].message.content)#.text.strip()
    except Exception as e:  # Catch a general exception
        print(f"An error occurred: {e}")
        return text  # Return the original text if an error occurs

# Function to process the PowerPoint file
def process_presentation(
    input_file_path: str, 
    output_file_path: str, 
    audience_type: str,
    executive_summary_slide = False):
    # Load the presentation
    prs = Presentation(input_file_path)
    
    # Iterate through each slide and each text box
    all_text = ""
    for slide_number, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                original_text = shape.text
                # Reword the text using GPT-4
                reworded_text = reword_text_with_gpt4(original_text, audience_type, shape.name)

                # Replace the original text with the reworded text
                shape.text = reworded_text
                all_text = all_text + ". " + original_text
        print(f"Processed slide {slide_number + 1}")
    
    if executive_summary_slide == True:
        print("Creating an executive summary slide")
        slide_layout = prs.slide_layouts[1]
        slide_exec = prs.slides.add_slide(slide_layout)
        slide_exec.placeholders[0].text = "Executive Summary"
        slide_exec.placeholders[1].text = create_exec_summary(all_text, audience_type)

        tf = slide_exec.placeholders[1].text_frame
            

        tf.word_wrap = True
        #for paragraph in tf.paragraphs:
        #    paragraph.font.size = Pt(26)

        feedback = create_ppt_feedback(all_text, audience_type)
        f = open("feedback.txt", "a")
        f.write(feedback)
        f.close()
    # Save the presentation
    prs.save(output_file_path)
    print(f"Presentation saved to {output_file_path}")



def getText(
    filename: str,
    file_format: str) -> str:
    if file_format == "docx":
        doc = Document(filename)
        fullText = []
        for para in doc.paragraphs:
            fullText.append(para.text)
        return '\n'.join(fullText)
    elif file_format =="pdf":
        reader = PdfReader(filename)
        text = ""
        for page in reader.pages:
            text += page.extract_text() + "\n"
        if len(text) > 5000:
            print("PDF file is too large so only taking the first 5,000 characters")
            text = text[:5000]
        return text

def document_to_ppt(
    doc: str,
    ppt: str, 
    audience_type: str, 
    slide_number: float, 
    file_format = "docx"):

    if file_format == "pptx":
        print("Rephrasing your powerpoint for "+ audience_type)
        process_presentation(
            input_file_path = doc, 
            output_file_path = ppt, 
            audience_type = audience_type,
            executive_summary_slide = True)
    else:
        text = getText(doc, file_format)

        try:
            # Use the OpenAI API to get a response
            response = openai.chat.completions.create(
                model="gpt-4-0613",  # Replace with the appropriate model
                messages=[
                    {
                        "role": "system", 
                        "content": (
                            "You are an assistant that provides powerpoints form word document inputs." +
                            " I am presenting to " + audience_type + ", so convert this document into a powerpoint suitable for this audience." +
                           "Ensure it has " + str(slide_number + 1) + " slides. The first slide is a title slide and the other slides are in the form of title and contents." +
                           "Provide  the answer as a dictionary for python of the form {'Slide 1':{'title':'title','subtitle: 'appropriate subtitle'},'Slide 2':{'title':'title','contents: unnumbered list of bullet points},'Slide 3':{'title':'title','contents: unnumbered list of bullet points}... etc}."+
                           " This string will be converted into a dictionary so ensure the bullet points are able to be read into a json as a list. The list length can't be longer than five. only print out the dictionary."

                                    )
                        },
                    {"role": "user", "content": text}
                ]
            )
            # Assuming the last message in the list will be the assistant's response
            ppt_dict =  str(response.choices[0].message.content)

            #print(ppt_dict)

            json_acceptable_string = ppt_dict.replace("'", "\"")

            #print(json_acceptable_string)
            try:
                ppt_dict = json.loads(json_acceptable_string, strict = False)

        
                prs = Presentation("theme.pptx")

                print("Making title slide")
                title_layout = prs.slide_layouts[0]
                normal_layout = prs.slide_layouts[1]
                title_slide = prs.slides.add_slide(title_layout)
                title_slide.placeholders[0].text = ppt_dict["Slide 1"]["title"]
                title_slide.placeholders[1].text = ppt_dict["Slide 1"]["subtitle"]

                for i in range(2,slide_number +2):
                    print("Making Slide Number " + str(i))
                    slide = prs.slides.add_slide(normal_layout)
                    slide.placeholders[0].text = ppt_dict["Slide "+str(i)]["title"]
                    tf = slide.placeholders[1].text_frame
            
                    for j in ppt_dict["Slide "+str(i)]["contents"]:
                 
                        p = tf.add_paragraph()
                        p.text = j
                        p.level = 0

                    tf.word_wrap = True
                    for paragraph in tf.paragraphs:
                        paragraph.font.size = Pt(26)
                # Save the presentation
                prs.save(ppt)
                print(f"Presentation saved to {ppt}")

            except Exception as e:  # Catch a general exception
                print(f"An error occurred with OpenAI Output- try rerunning the query")

        except Exception as e:  # Catch a general exception
            print(f"An error occurred: {e}")

    
# Example usage
wd = "C://Users//Administrator//Documents//GitHub//genai_group1"
import os
os.chdir(wd)
input_file_path =  "inputs//GPTB4.pptx"
output_file_path = "outputs//GPTA4"
output_file_path_technical = "outputs//GPTA4_tech.pptx"
output_file_path_baby = "outputs//GPTA4_babies.pptx"
audience_type1 = "a technical audience"
audience_type2 = "a bunch of five year olds who like thomas the tank engine"
audience_type3 = "people who know nothing about finance and only knowledge of the economy is robux gift cards and use that as a point of reference"

#Worked examples

document_to_ppt(
    doc = input_file_path, 
    ppt = output_file_path_technical,
    audience_type = audience_type1, 
    slide_number = 3, 
    file_format = "pptx"
    )
#document_to_ppt(
#    doc = "inputs//Singapore Day 1.docx",
#    ppt = "outputs//singapore.pptx", 
#    audience_type = audience_type2, 
#    slide_number = 3, 
#    file_format = "docx"
#    )

#document_to_ppt(
#    doc = "inputs//What is Sharpe Ratio.docx",
#    ppt = "outputs//SSharpe.pptx",
#    audience_type = audience_type2, 
#    slide_number = 2, 
#    file_format = "docx"
#    )

#document_to_ppt(
#    doc = "inputs//What is Sharpe Ratio.docx",
#    ppt = "outputs//SSharpe2.pptx",
#    audience_type = audience_type3, 
#    slide_number = 5, 
#    file_format = "docx"
#    )

#document_to_ppt(
#    doc = "inputs//abm_cc.pdf",
#    ppt = "outputs//abm_robux.pptx",
#    audience_type = audience_type3, 
#    slide_number = 5, 
#    file_format = "pdf"
#    )
#process_presentation(
#    input_file_path, 
#    output_file_path_technical, 
#    audience_type1,
#    executive_summary_slide = True)

#process_presentation(
#    input_file_path, 
#    output_file_path_baby, 
#    audience_type2)
